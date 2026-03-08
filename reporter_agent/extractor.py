from __future__ import annotations

import re
from collections import Counter
from pathlib import Path

from .models import SlideRecord


SECTION_PATTERNS: dict[str, list[str]] = {
    "objective": [r"\bobjective\b", r"\baim\b", r"\bgoal\b", r"\bscope\b"],
    "methodology": [r"\bmethod", r"\bapproach\b", r"\bsetup\b", r"\bassumption"],
    "calibration": [r"\bcalibration\b", r"\btuning\b", r"\bfit\b"],
    "sensitivity": [r"\bsensitivity\b", r"\bparameter sweep\b", r"\bvariation\b"],
    "results": [r"\bresult", r"\boutput\b", r"\bperformance\b", r"\bcomparison\b"],
    "conclusion": [r"\bconclusion\b", r"\bsummary\b", r"\brecommendation\b", r"\bnext step\b"],
    "literature_review": [r"\bliterature\b", r"\bpaper\b", r"\brelated work\b", r"\breview\b"],
}

STOPWORDS = {
    "the",
    "a",
    "an",
    "and",
    "or",
    "to",
    "for",
    "of",
    "in",
    "on",
    "with",
    "from",
    "by",
    "is",
    "are",
    "was",
    "were",
    "be",
    "this",
    "that",
    "it",
    "as",
    "at",
    "we",
    "our",
}


def extract_text_from_slide(slide) -> tuple[str, str]:
    title = ""
    lines: list[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = (shape.text or "").strip()
        if not text:
            continue
        if shape == slide.shapes.title and not title:
            title = text
        lines.append(text)
    raw_text = "\n".join(lines)
    return title, raw_text


def detect_section(title: str, raw_text: str) -> str:
    haystack = f"{title}\n{raw_text}".lower()
    scores = {section: 0 for section in SECTION_PATTERNS}
    for section, patterns in SECTION_PATTERNS.items():
        for pattern in patterns:
            if re.search(pattern, haystack):
                scores[section] += 1
    best_section = max(scores, key=scores.get)
    return best_section if scores[best_section] > 0 else "general"


def extract_key_phrases(text: str, top_k: int = 8) -> list[str]:
    tokens = re.findall(r"[a-zA-Z][a-zA-Z0-9_\-]{2,}", text.lower())
    filtered = [t for t in tokens if t not in STOPWORDS]
    counts = Counter(filtered)
    return [token for token, _ in counts.most_common(top_k)]


def extract_slide_records(pptx_path: Path) -> list[SlideRecord]:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "python-pptx is required for PPT ingestion. Install with: python -m pip install python-pptx"
        ) from exc

    prs = Presentation(str(pptx_path))
    slide_records: list[SlideRecord] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title, raw_text = extract_text_from_slide(slide)
        section = detect_section(title, raw_text)
        key_phrases = extract_key_phrases(f"{title}\n{raw_text}")
        slide_records.append(
            SlideRecord(
                source_file=str(pptx_path),
                slide_index=idx,
                title=title or f"Slide {idx}",
                raw_text=raw_text,
                section=section,
                key_phrases=key_phrases,
            )
        )
    return slide_records
