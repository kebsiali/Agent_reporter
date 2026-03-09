from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from .models import ReportPlan


def export_plan_json(plan: ReportPlan, output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(plan.to_dict(), indent=2), encoding="utf-8")


def export_plan_markdown(plan: ReportPlan, output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    lines: list[str] = []
    lines.append(f"# Report Plan: {plan.task_name}")
    lines.append("")
    lines.append(f"- Report type: `{plan.report_type}`")
    lines.append(f"- Generated at: `{plan.created_at}`")
    lines.append("")
    lines.append("## Assumptions")
    for a in plan.assumptions:
        lines.append(f"- {a}")
    lines.append("")

    for s in plan.slides:
        lines.append(f"## Slide {s.slide_number}: {s.title} ({s.section})")
        lines.append(f"Objective: {s.objective}")
        lines.append(f"Confidence: `{s.confidence_label}` ({s.confidence:.3f})")
        lines.append("")
        lines.append("Auto-fill draft:")
        lines.append("")
        lines.append(s.autofill_text)
        lines.append("")
        lines.append("Evidence gaps:")
        if s.evidence_gaps:
            for gap in s.evidence_gaps:
                lines.append(f"- {gap}")
        else:
            lines.append("- None")
        lines.append("")
        lines.append("Placeholders:")
        for p in s.placeholders:
            lines.append(f"- {p}")
        lines.append("")
        lines.append("Missing info guidance:")
        for g in s.missing_info_guidance:
            lines.append(f"- {g}")
        lines.append("")
        lines.append("Source examples:")
        if s.source_examples:
            for src in s.source_examples:
                lines.append(f"- {src}")
        else:
            lines.append("- None")
        lines.append("")

    output_path.write_text("\n".join(lines), encoding="utf-8")


def _apply_profile_to_paragraph(paragraph, font_name: str | None, font_size_pt: float | None):
    if font_name:
        paragraph.font.name = font_name
    if font_size_pt:
        from pptx.util import Pt

        paragraph.font.size = Pt(font_size_pt)


def _find_main_content_text_frame(slide, prs):
    # Prefer explicit body/content placeholders with largest area.
    candidates = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape == slide.shapes.title:
            continue
        if getattr(shape, "is_placeholder", False):
            area = int(shape.width) * int(shape.height)
            candidates.append((area, shape.text_frame))

    if candidates:
        candidates.sort(key=lambda x: x[0], reverse=True)
        return candidates[0][1]

    # Fallback: create a main body textbox inside slide canvas.
    from pptx.util import Inches

    left = Inches(0.7)
    top = Inches(1.5)
    width = prs.slide_width - Inches(1.4)
    height = prs.slide_height - Inches(2.0)
    box = slide.shapes.add_textbox(left, top, width, height)
    return box.text_frame


def export_plan_pptx(
    plan: ReportPlan,
    output_path: Path,
    template_pptx: Path | None = None,
    style_profile: dict[str, Any] | None = None,
) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "python-pptx is required for PPT export. Install with: python -m pip install python-pptx"
        ) from exc

    output_path.parent.mkdir(parents=True, exist_ok=True)
    if template_pptx and template_pptx.exists():
        prs = Presentation(str(template_pptx))
    else:
        prs = Presentation()

    title_font_name = style_profile.get("title_font_name") if style_profile else None
    title_font_size = style_profile.get("title_font_size_pt") if style_profile else None
    body_font_name = style_profile.get("body_font_name") if style_profile else None
    body_font_size = style_profile.get("body_font_size_pt") if style_profile else None

    for planned in plan.slides:
        layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = f"{planned.slide_number}. {planned.title}"
            if slide.shapes.title.text_frame and slide.shapes.title.text_frame.paragraphs:
                _apply_profile_to_paragraph(
                    slide.shapes.title.text_frame.paragraphs[0], title_font_name, title_font_size
                )

        body = _find_main_content_text_frame(slide, prs)
        body.clear()

        p = body.paragraphs[0]
        p.text = f"Section: {planned.section}"
        p.level = 0
        _apply_profile_to_paragraph(p, body_font_name, body_font_size)

        p = body.add_paragraph()
        p.text = f"Confidence: {planned.confidence_label} ({planned.confidence:.3f})"
        p.level = 0
        _apply_profile_to_paragraph(p, body_font_name, body_font_size)

        p = body.add_paragraph()
        p.text = "Draft:"
        p.level = 0
        _apply_profile_to_paragraph(p, body_font_name, body_font_size)

        p = body.add_paragraph()
        p.text = planned.autofill_text[:1200]
        p.level = 1
        _apply_profile_to_paragraph(p, body_font_name, body_font_size)

        p = body.add_paragraph()
        p.text = "Fill these:"
        p.level = 0
        _apply_profile_to_paragraph(p, body_font_name, body_font_size)
        for placeholder in planned.placeholders:
            p = body.add_paragraph()
            p.text = placeholder
            p.level = 1
            _apply_profile_to_paragraph(p, body_font_name, body_font_size)

        p = body.add_paragraph()
        p.text = "Evidence gaps:"
        p.level = 0
        _apply_profile_to_paragraph(p, body_font_name, body_font_size)
        if planned.evidence_gaps:
            for gap in planned.evidence_gaps:
                p = body.add_paragraph()
                p.text = gap
                p.level = 1
                _apply_profile_to_paragraph(p, body_font_name, body_font_size)
        else:
            p = body.add_paragraph()
            p.text = "None"
            p.level = 1
            _apply_profile_to_paragraph(p, body_font_name, body_font_size)

        p = body.add_paragraph()
        p.text = "Missing info guidance:"
        p.level = 0
        _apply_profile_to_paragraph(p, body_font_name, body_font_size)
        for g in planned.missing_info_guidance:
            p = body.add_paragraph()
            p.text = g
            p.level = 1
            _apply_profile_to_paragraph(p, body_font_name, body_font_size)

    prs.save(str(output_path))
