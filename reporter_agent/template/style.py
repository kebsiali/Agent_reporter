from __future__ import annotations

import json
from pathlib import Path
from typing import Any


def _font_size_points(font) -> float | None:
    try:
        return float(font.size.pt) if font.size else None
    except Exception:  # noqa: BLE001
        return None


def extract_template_profile(template_pptx: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise ModuleNotFoundError(
            "python-pptx is required for template profiling. Install with: python -m pip install python-pptx"
        ) from exc

    prs = Presentation(str(template_pptx))
    profile: dict[str, Any] = {
        "template_path": str(template_pptx),
        "title_font_name": None,
        "title_font_size_pt": None,
        "body_font_name": None,
        "body_font_size_pt": None,
    }

    for slide in prs.slides:
        if slide.shapes.title and hasattr(slide.shapes.title, "text_frame"):
            tf = slide.shapes.title.text_frame
            if tf.paragraphs:
                p = tf.paragraphs[0]
                if p.runs:
                    run = p.runs[0]
                    if profile["title_font_name"] is None:
                        profile["title_font_name"] = run.font.name
                    if profile["title_font_size_pt"] is None:
                        profile["title_font_size_pt"] = _font_size_points(run.font)
        for shape in slide.shapes:
            if shape == slide.shapes.title or not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                run = tf.paragraphs[0].runs[0]
                if profile["body_font_name"] is None:
                    profile["body_font_name"] = run.font.name
                if profile["body_font_size_pt"] is None:
                    profile["body_font_size_pt"] = _font_size_points(run.font)
        if all(profile.get(k) is not None for k in ["title_font_name", "title_font_size_pt", "body_font_name", "body_font_size_pt"]):
            break

    return profile


def save_template_profile(path: Path, profile: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(profile, indent=2), encoding="utf-8")


def load_template_profile(path: Path) -> dict[str, Any] | None:
    if not path.exists():
        return None
    return json.loads(path.read_text(encoding="utf-8"))

